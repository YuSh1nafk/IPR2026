from pptx import Presentation
from pptx.util import Inches, Pt
import os
import base64
import io
import re


def create_unified_ppt(all_slide_decks, image_cache, output_filename="outputs/Full_Unit_Slides.pptx",
                       template_path="templates/template.pptx"):
    prs = Presentation(template_path) if os.path.exists(template_path) else Presentation()

    TITLE_LAYOUT = prs.slide_layouts[0]
    CONTENT_LAYOUT = prs.slide_layouts[1]

    for deck in all_slide_decks:

        slide_title = prs.slides.add_slide(TITLE_LAYOUT)

        if slide_title.shapes.title:
            slide_title.shapes.title.text = deck.unit_title

        for shape in slide_title.placeholders:
            if shape != slide_title.shapes.title and shape.has_text_frame:
                shape.text = deck.lesson_name
                break

        for slide_info in deck.slides:
            slide = prs.slides.add_slide(CONTENT_LAYOUT)

            if slide.shapes.title:
                slide.shapes.title.text = slide_info.title

            body_shape = None
            for shape in slide.placeholders:
                if shape != slide.shapes.title and shape.has_text_frame:
                    body_shape = shape
                    break

            img_left = Inches(1)
            img_top = Inches(2)
            if body_shape:
                tf = body_shape.text_frame
                tf.clear()

                for line in slide_info.content_lines:
                    match = re.search(r'\[IMAGE_TAG_(\d+)\]', line)
                    if match:
                        seq_id = match.group(1)
                        if seq_id in image_cache:
                            try:
                                img_data = base64.b64decode(image_cache[seq_id])
                                img_stream = io.BytesIO(img_data)
                                slide.shapes.add_picture(img_stream, left=img_left, top=img_top, width=Inches(4.5))
                                img_top += Inches(3)
                            except Exception as e:
                                print(f"Lỗi chèn ảnh {seq_id}: {e}")
                    else:
                        p = tf.add_paragraph()
                        p.text = line
                        p.level = 0
            if slide_info.notes and slide.has_notes_slide:
                slide.notes_slide.notes_text_frame.text = slide_info.notes

    os.makedirs(os.path.dirname(output_filename), exist_ok=True)
    prs.save(output_filename)
    print(f"🚀 Xuất Slide thành công! Đã chèn cả hình ảnh. File tại: {output_filename}")